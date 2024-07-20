import os
import random
import re
from transformers import pipeline
from pptx import Presentation

# Initialize the model and pipeline
generator = pipeline('text-generation', model='gpt2')

# Function to generate PPT text using the Hugging Face model
def create_ppt_text(input_string):
    prompt = (
        "Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.\n"
        "Notice:\n"
        "- You do all the presentation text for the user.\n"
        "- You write the texts no longer than 250 characters!\n"
        "- You make very short titles!\n"
        "- You make the presentation easy to understand.\n"
        "- The presentation has a table of contents.\n"
        "- The presentation has a summary.\n"
        "- At least 7 slides.\n\n"
        "Example! - Stick to this formatting exactly!\n"
        "#Title: TITLE OF THE PRESENTATION\n\n"
        "#Slide: 1\n"
        "#Header: table of contents\n"
        "#Content: 1. CONTENT OF THIS POWERPOINT\n"
        "2. CONTENTS OF THIS POWERPOINT\n"
        "3. CONTENT OF THIS POWERPOINT\n...\n\n"
        "#Slide: 2\n"
        "#Header: TITLE OF SLIDE\n"
        "#Content: CONTENT OF THE SLIDE\n\n"
        "#Slide: 3\n"
        "#Header: TITLE OF SLIDE\n"
        "#Content: CONTENT OF THE SLIDE\n\n"
        "#Slide: 4\n"
        "#Header: TITLE OF SLIDE\n"
        "#Content: CONTENT OF THE SLIDE\n\n"
        "#Slide: 5\n"
        "#Header: summary\n"
        "#Content: CONTENT OF THE SUMMARY\n\n"
        "#Slide: END\n"
        "The user wants a presentation about " + input_string
    )
    
    response = generator(prompt, max_length=2048, num_return_sequences=1)
    ppt_text = response[0]['generated_text'].strip()
    return ppt_text

# Function to create a PowerPoint presentation
def create_ppt(text_content, design_number, ppt_name):
    prs = Presentation(f"Powerpointer/Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True

    lines = text_content.split('\n')
    for line in lines:
        if line.startswith('#Title:'):
            header = line.replace('#Title:', '').strip()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = header
            body_shape = slide.shapes.placeholders[1]
            continue
        elif line.startswith('#Slide:'):
            if slide_count > 0:
                slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[slide_placeholder_index]
                tf = body_shape.text_frame
                tf.text = content
            content = "" 
            slide_count += 1
            slide_layout_index = last_slide_layout_index
            layout_indices = [1, 7, 8] 
            while slide_layout_index == last_slide_layout_index:
                if firsttime == True:
                    slide_layout_index = 1
                    slide_placeholder_index = 1
                    firsttime = False
                    break
                slide_layout_index = random.choice(layout_indices)
                if slide_layout_index == 8:
                    slide_placeholder_index = 2
                else:
                    slide_placeholder_index = 1
            last_slide_layout_index = slide_layout_index
            continue

        elif line.startswith('#Header:'):
            header = line.replace('#Header:', '').strip()
            continue

        elif line.startswith('#Content:'):
            content = line.replace('#Content:', '').strip()
            continue

    prs.save(f'Powerpointer/GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"Powerpointer/GeneratedPresentations/{ppt_name}.pptx"
    
    return file_path

# Main function to get the bot response and create the PowerPoint
def get_bot_response(msg):
    user_text = msg
    last_char = user_text[-1]
    input_string = user_text
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', input_string)
    input_string = input_string.replace("\n", "")
    number = 1

    if last_char.isdigit():
        number = int(last_char)
        input_string = user_text[:-2]
        print("Design Number:", number, "selected.")
    else:
        print("No design specified, using default design...")
    
    if number > 5:  # assuming we have 5 designs
        number = 1

    ppt_text = create_ppt_text(input_string)
    ppt_file_path = create_ppt(ppt_text, number, input_string)

    return ppt_file_path

# Example usage
if __name__ == "__main__":
    topic = "AI Revolution"
    print(get_bot_response(topic))
